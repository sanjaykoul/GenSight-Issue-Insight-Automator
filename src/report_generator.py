
# -*- coding: utf-8 -*-
# src/report_generator.py
"""
Report generators (PDF & PPTX) for GenSight-Issue-Insight-Automator.

Features
--------
- Auto-creates `reports/<MONTH>/` and `reports/<MONTH>/charts/`
- PDF: Title + key metrics + optional AI text + charts (multi-page)
- PPTX: Title, key highlights, optional AI insights slide, chart slides
- Optional Month-over-Month (MoM) bar chart generation and slide/page
- Optional logo on title slide / PDF cover

Dependencies
------------
pip install reportlab python-pptx matplotlib
"""

import os
from typing import List, Dict, Optional

# PDF
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader

# PPTX
from pptx import Presentation
from pptx.util import Inches, Pt

# For optional MoM chart image
import matplotlib.pyplot as plt


# ----------------------------- #
#            PATHS              #
# ----------------------------- #

def ensure_month_folder(month_label: str) -> str:
    """
    Ensure reports/<MONTH>/ and /charts exist. Return month folder path.
    """
    folder = f"reports/{month_label}"
    os.makedirs(folder, exist_ok=True)
    os.makedirs(os.path.join(folder, "charts"), exist_ok=True)
    return folder


# ----------------------------- #
#       UTILS (PDF TEXT)        #
# ----------------------------- #

def _wrap_text(text: str, max_chars: int = 100) -> List[str]:
    """
    Simple word-wrap for AI text paragraphs.
    """
    if not text:
        return []
    words = str(text).split()
    lines, cur = [], []
    count = 0
    for w in words:
        if count + len(w) + (1 if cur else 0) > max_chars:
            lines.append(" ".join(cur))
            cur, count = [w], len(w)
        else:
            cur.append(w)
            count += len(w) + (1 if cur else 0)
    if cur:
        lines.append(" ".join(cur))
    return lines


# ----------------------------- #
#     OPTIONAL: MoM CHART       #
# ----------------------------- #

def build_mom_chart_image(summary: Dict, month_label_for_folder: str) -> Optional[str]:
    """
    Create a Month-over-Month (MoM) bar chart image from summary['by_month'].
    Returns path to the saved image or None.
    """
    by_month_df = summary.get("by_month")
    if by_month_df is None or by_month_df.empty:
        return None

    # Sort by month label for a clean x-axis
    df = by_month_df.sort_values("month_label").reset_index(drop=True)
    labels = df["month_label"].tolist()
    values = df["issue_count"].tolist()

    if not labels:
        return None

    folder = ensure_month_folder(month_label_for_folder)
    charts_dir = os.path.join(folder, "charts")
    out_path = os.path.join(charts_dir, "mom_comparison.png")

    plt.figure(figsize=(7.5, 4.5))
    plt.bar(labels, values, color="#6C5CE7")
    plt.title("Month-over-Month Issue Volume")
    plt.xlabel("Month")
    plt.ylabel("Issues")
    plt.xticks(rotation=20)
    plt.tight_layout()
    plt.savefig(out_path, dpi=150)
    plt.close()
    return out_path


# ----------------------------- #
#           PDF MAKER           #
# ----------------------------- #

def _pdf_cover_page(c: canvas.Canvas,
                    summary: Dict,
                    month_label: str,
                    ai_text: Optional[str] = None,
                    logo_path: Optional[str] = None):
    """
    Draw a cover page: title, key metrics, optional logo and AI text block.
    """
    width, height = A4
    margin_left = 40
    y = height - 60

    # Logo (optional, top-right)
    if logo_path and os.path.exists(logo_path):
        try:
            img = ImageReader(logo_path)
            # Keep a modest size to avoid layout issues
            c.drawImage(img, width - 150, height - 110, width=100, height=60, preserveAspectRatio=True, anchor="ne")
        except Exception:
            pass

    # Title
    c.setFont("Helvetica-Bold", 18)
    c.drawString(margin_left, y, f"Monthly Issue Report — {month_label}")

    # Key metrics
    y -= 30
    c.setFont("Helvetica", 12)
    by_status = summary.get("by_status", {})
    by_issue = summary.get("by_issue_type", {})
    by_engineer = summary.get("by_engineer", {})
    top_eng = list(by_engineer.items())[:3]

    c.drawString(margin_left, y, f"Status: {by_status}")
    y -= 18
    c.drawString(margin_left, y, f"Issue Types: {by_issue}")
    y -= 18
    c.drawString(margin_left, y, f"Top Engineers: {top_eng}")

    # AI text block (optional)
    if ai_text:
        y -= 28
        c.setFont("Helvetica-Bold", 12)
        c.drawString(margin_left, y, "AI Summary:")
        y -= 16
        c.setFont("Helvetica", 11)
        wrapped = _wrap_text(ai_text, max_chars=100)
        for line in wrapped:
            if y < 60:  # new page if near bottom
                c.showPage()
                y = height - 60
                c.setFont("Helvetica", 11)
            c.drawString(margin_left, y, line)
            y -= 14

    c.showPage()


def _pdf_charts_pages(c: canvas.Canvas, chart_paths: List[str]):
    """
    Put up to two charts per page, auto-paginate.
    """
    if not chart_paths:
        return

    width, height = A4
    margin_left = 40
    y_img = height - 120
    charts_on_page = 0

    c.setFont("Helvetica-Bold", 16)
    c.drawString(margin_left, height - 60, "Charts")

    for p in chart_paths:
        try:
            img = ImageReader(p)
            c.drawImage(img, margin_left, y_img - 220, width=520, height=200,
                        preserveAspectRatio=True, anchor="c")
            y_img -= 240
            charts_on_page += 1
            if charts_on_page == 2:
                c.showPage()
                c.setFont("Helvetica-Bold", 16)
                c.drawString(margin_left, height - 60, "Charts (contd.)")
                y_img = height - 120
                charts_on_page = 0
        except Exception:
            # skip unreadable images
            continue

    c.showPage()


def generate_pdf_report(summary: Dict,
                        month_label: str,
                        charts: Optional[List[str]] = None,
                        ai_text: Optional[str] = None,
                        logo_path: Optional[str] = None,
                        include_mom: bool = True,
                        output_path: Optional[str] = None) -> str:
    """
    Create a PDF report for a given month.
    Optionally include: AI text block, logo on cover, and a MoM chart page.
    """
    folder = ensure_month_folder(month_label)
    if output_path is None:
        output_path = os.path.join(folder, f"Monthly_Report_{month_label}.pdf")

    c = canvas.Canvas(output_path, pagesize=A4)

    # Cover page with metrics + AI text + logo
    _pdf_cover_page(c, summary, month_label, ai_text=ai_text, logo_path=logo_path)

    # Charts pages (user-provided charts)
    chart_paths = charts[:] if charts else []

    # Optional MoM chart page
    if include_mom:
        mom_img = build_mom_chart_image(summary, month_label_for_folder=month_label)
        if mom_img:
            chart_paths.append(mom_img)

    _pdf_charts_pages(c, chart_paths)

    # Save
    c.save()
    return output_path


# ----------------------------- #
#           PPT MAKER           #
# ----------------------------- #

def _ppt_title_slide(prs: Presentation, month_label: str, logo_path: Optional[str] = None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"Monthly Issue Insights — {month_label}"
    subtitle = slide.placeholders[1]
    subtitle.text = "Auto-generated by GenSight-Issue-Insight-Automator"

    # Optional logo (top-right)
    if logo_path and os.path.exists(logo_path):
        try:
            left = Inches(9.0)
            top = Inches(0.2)
            height = Inches(0.8)
            slide.shapes.add_picture(logo_path, left, top, height=height)
        except Exception:
            pass


def _ppt_highlights_slide(prs: Presentation, summary: Dict):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Highlights"

    tf = slide.placeholders[1].text_frame
    tf.clear()

    by_status = summary.get("by_status", {})
    by_issue = summary.get("by_issue_type", {})
    by_engineer = summary.get("by_engineer", {})
    top_eng = list(by_engineer.items())[:3]

    p = tf.paragraphs[0]
    p.text = f"Status: {by_status}"
    p.level = 0

    p = tf.add_paragraph()
    p.text = f"Issue Types: {by_issue}"
    p.level = 0

    p = tf.add_paragraph()
    p.text = f"Top Engineers: {top_eng}"
    p.level = 0


def _ppt_ai_slide(prs: Presentation, ai_text: str):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "AI Summary"

    tf = slide.placeholders[1].text_frame
    tf.clear()

    # Wrap gently to avoid oversized bullets
    wrapped = _wrap_text(ai_text, max_chars=110) if ai_text else []
    if not wrapped:
        wrapped = ["No AI summary available."]

    # First paragraph
    tf.paragraphs[0].text = wrapped[0]
    tf.paragraphs[0].level = 0
    for line in wrapped[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0


def _ppt_chart_slide(prs: Presentation, image_path: str, title: Optional[str] = None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = title or os.path.basename(image_path)
    left, top, height = Inches(1), Inches(1.5), Inches(4.5)
    try:
        slide.shapes.add_picture(image_path, left, top, height=height)
    except Exception:
        # ignore missing image
        pass


def _ppt_mom_slide(prs: Presentation, summary: Dict, mom_chart_path: Optional[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Month-over-Month Comparison"

    tf = slide.placeholders[1].text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Issue volume by month:"
    p.level = 0

    by_month_df = summary.get("by_month")
    if by_month_df is not None and not by_month_df.empty:
        for _, row in by_month_df.sort_values("month_label").iterrows():
            pp = tf.add_paragraph()
            pp.text = f"{row['month_label']}: {row['issue_count']} issues"
            pp.level = 1

    # If we also built a chart image, add a visual slide too
    if mom_chart_path and os.path.exists(mom_chart_path):
        _ppt_chart_slide(prs, mom_chart_path, title="MoM Issue Volume")


def generate_ppt_report(summary: Dict,
                        month_label: str,
                        charts: Optional[List[str]] = None,
                        ai_text: Optional[str] = None,
                        logo_path: Optional[str] = None,
                        include_mom: bool = True,
                        output_path: Optional[str] = None) -> str:
    """
    Create a PPTX for a given month, with:
      - Title slide (+optional logo)
      - Key highlights
      - Optional AI summary slide
      - Chart slides
      - Optional MoM summary + chart slide
    """
    folder = ensure_month_folder(month_label)
    if output_path is None:
        output_path = os.path.join(folder, f"Monthly_Issue_Insights_{month_label}.pptx")

    prs = Presentation()

    # Title + highlights
    _ppt_title_slide(prs, month_label, logo_path=logo_path)
    _ppt_highlights_slide(prs, summary)

    # Optional AI slide
    if ai_text:
        _ppt_ai_slide(prs, ai_text)

    # Chart slides
    if charts:
        for img in charts:
            if img and os.path.exists(img):
                _ppt_chart_slide(prs, img)

    # Optional MoM (text + optional chart)
    mom_chart_path = None
    if include_mom:
        mom_chart_path = build_mom_chart_image(summary, month_label_for_folder=month_label)
        _ppt_mom_slide(prs, summary, mom_chart_path)

    prs.save(output_path)
    return output_path